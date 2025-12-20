"""
OCR Service using Groq API for text extraction from images and videos
"""
import logging
import base64
import os
import tempfile
import subprocess
from typing import Optional, List
from groq import Groq
from app.config import settings
import cv2
import numpy as np
from PIL import Image
import io

logger = logging.getLogger(__name__)

class OCRService:
    def __init__(self):
        self.client = Groq(api_key=settings.GROQ_API_KEY)
        self.model = "meta-llama/llama-4-scout-17b-16e-instruct"  # Using Llama Vision for OCR
    
    def _encode_image_to_base64(self, image_path: str) -> str:
        """Convert image to base64 string"""
        try:
            with open(image_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        except Exception as e:
            logger.error(f"Error encoding image to base64: {e}")
            raise
    
    def _find_document_contour(self, img):
        """Find the largest quadrilateral contour (document boundary)"""
        try:
            # Convert to grayscale if needed
            if len(img.shape) == 3:
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            else:
                gray = img
            
            # Apply Gaussian blur to reduce noise
            blurred = cv2.GaussianBlur(gray, (5, 5), 0)
            
            # Edge detection
            edged = cv2.Canny(blurred, 75, 200)
            
            # Find contours
            contours, _ = cv2.findContours(edged.copy(), cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
            contours = sorted(contours, key=cv2.contourArea, reverse=True)[:5]
            
            # Find the largest contour that is a quadrilateral
            for contour in contours:
                peri = cv2.arcLength(contour, True)
                approx = cv2.approxPolyDP(contour, 0.02 * peri, True)
                
                if len(approx) == 4:
                    return approx
            
            return None
        except Exception as e:
            logger.warning(f"Document contour detection failed: {e}")
            return None
    
    def _order_points(self, pts):
        """Order points in top-left, top-right, bottom-right, bottom-left order"""
        rect = np.zeros((4, 2), dtype="float32")
        
        # Sum and difference to find corners
        s = pts.sum(axis=1)
        diff = np.diff(pts, axis=1)
        
        rect[0] = pts[np.argmin(s)]      # Top-left
        rect[2] = pts[np.argmax(s)]      # Bottom-right
        rect[1] = pts[np.argmin(diff)]   # Top-right
        rect[3] = pts[np.argmax(diff)]   # Bottom-left
        
        return rect
    
    def _perspective_transform(self, img, contour):
        """Apply perspective transform to get bird's eye view of document"""
        try:
            # Reshape contour and order points
            pts = contour.reshape(4, 2)
            rect = self._order_points(pts)
            
            (tl, tr, br, bl) = rect
            
            # Compute width of new image
            widthA = np.sqrt(((br[0] - bl[0]) ** 2) + ((br[1] - bl[1]) ** 2))
            widthB = np.sqrt(((tr[0] - tl[0]) ** 2) + ((tr[1] - tl[1]) ** 2))
            maxWidth = max(int(widthA), int(widthB))
            
            # Compute height of new image
            heightA = np.sqrt(((tr[0] - br[0]) ** 2) + ((tr[1] - br[1]) ** 2))
            heightB = np.sqrt(((tl[0] - bl[0]) ** 2) + ((tl[1] - bl[1]) ** 2))
            maxHeight = max(int(heightA), int(heightB))
            
            # Destination points for perspective transform
            dst = np.array([
                [0, 0],
                [maxWidth - 1, 0],
                [maxWidth - 1, maxHeight - 1],
                [0, maxHeight - 1]
            ], dtype="float32")
            
            # Calculate perspective transform matrix and apply it
            M = cv2.getPerspectiveTransform(rect, dst)
            warped = cv2.warpPerspective(img, M, (maxWidth, maxHeight))
            
            return warped
        except Exception as e:
            logger.warning(f"Perspective transform failed: {e}")
            return img
    
    def _remove_shadows(self, img):
        """Remove shadows from image for better text contrast"""
        try:
            # Convert to LAB color space
            lab = cv2.cvtColor(img, cv2.COLOR_BGR2LAB)
            l, a, b = cv2.split(lab)
            
            # Apply CLAHE to L channel
            clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8, 8))
            l = clahe.apply(l)
            
            # Merge channels
            lab = cv2.merge([l, a, b])
            
            # Convert back to BGR
            result = cv2.cvtColor(lab, cv2.COLOR_LAB2BGR)
            return result
        except Exception as e:
            logger.warning(f"Shadow removal failed: {e}")
            return img
    
    def _enhance_contrast(self, img):
        """Enhance image contrast using CLAHE"""
        try:
            if len(img.shape) == 3:
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            else:
                gray = img
            
            # Apply CLAHE
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            enhanced = clahe.apply(gray)
            
            return enhanced
        except Exception as e:
            logger.warning(f"Contrast enhancement failed: {e}")
            return img if len(img.shape) == 2 else cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    
    def _smart_binarize(self, img):
        """Apply intelligent binarization for better text extraction"""
        try:
            # Ensure grayscale
            if len(img.shape) == 3:
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            else:
                gray = img
            
            # Apply denoising
            denoised = cv2.fastNlMeansDenoising(gray, h=10)
            
            # Apply Otsu's thresholding
            _, binary = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            return binary
        except Exception as e:
            logger.warning(f"Binarization failed: {e}")
            return img
    
    def _preprocess_image_for_ocr(self, image_path: str) -> str:
        """Advanced preprocessing with document scanning capabilities"""
        try:
            # Read image
            img = cv2.imread(image_path)
            if img is None:
                raise ValueError("Could not read image file")
            
            original_img = img.copy()
            
            # Step 1: Try to find and correct document perspective
            contour = self._find_document_contour(img)
            if contour is not None:
                logger.info("Document boundary detected, applying perspective correction")
                img = self._perspective_transform(img, contour)
            
            # Step 2: Remove shadows for better uniformity
            img = self._remove_shadows(img)
            
            # Step 3: Enhance contrast
            enhanced = self._enhance_contrast(img)
            
            # Step 4: Apply smart binarization
            binary = self._smart_binarize(enhanced)
            
            # Step 5: Apply sharpening for better text clarity
            kernel = np.array([[-1,-1,-1], [-1,9,-1], [-1,-1,-1]])
            sharpened = cv2.filter2D(binary, -1, kernel)
            
            # Save preprocessed image to temporary file
            temp_path = image_path.replace('.', '_processed.')
            cv2.imwrite(temp_path, sharpened)
            
            logger.info("Successfully preprocessed image with document scanning pipeline")
            return temp_path
            
        except Exception as e:
            logger.warning(f"Advanced preprocessing failed, using original: {e}")
            return image_path
    
    async def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using Groq Vision API"""
        try:
            # Preprocess image for better OCR
            processed_image_path = self._preprocess_image_for_ocr(image_path)
            
            # Encode image to base64
            base64_image = self._encode_image_to_base64(processed_image_path)
            
            # Create the prompt for OCR
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": """Please extract all text content from this image. 
                            Focus on:
                            1. All readable text, including handwritten text if present
                            2. Maintain the structure and formatting as much as possible
                            3. Include any numbers, formulas, or equations
                            4. If there are tables, preserve the table structure
                            5. If no text is found, respond with 'No readable text found in image'
                            
                            Please provide only the extracted text without any additional commentary."""
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ]
            
            # Call Groq API
            response = self.client.chat.completions.create(
                model=self.model,
                messages=messages,
                max_tokens=2000,
                temperature=0.1
            )
            
            extracted_text = response.choices[0].message.content.strip()
            
            # Clean up preprocessed image if it was created
            if processed_image_path != image_path and os.path.exists(processed_image_path):
                os.remove(processed_image_path)
            
            logger.info(f"Successfully extracted text from image: {len(extracted_text)} characters")
            
            # Apply AI formatting to clean up the text
            formatted_text = await self._format_ocr_text_with_ai(extracted_text)
            return formatted_text
            
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return f"OCR extraction failed: {str(e)}"
    
    async def _format_ocr_text_with_ai(self, raw_text: str) -> str:
        """Use AI to clean up and format OCR text"""
        try:
            # Skip formatting if text is too short or already well-formatted
            if len(raw_text.strip()) < 50 or "No readable text found" in raw_text:
                return raw_text
            
            messages = [
                {
                    "role": "system",
                    "content": """You are a text formatting assistant. Your job is to clean up OCR-extracted text and make it readable with proper spacing and structure.

Rules:
1. Remove duplicate or redundant text
2. Fix obvious OCR errors and typos
3. Organize content with clear hierarchical structure
4. Use proper line breaks and spacing:
   - Double line break between major sections
   - Single line break between related items
   - Indent sub-points with spaces or dashes
5. Keep all original information - don't remove content
6. DON'T use markdown symbols (no ##, **, -, etc.)
7. Use plain text formatting with line breaks and indentation
8. Make headings stand out with CAPITAL LETTERS or by spacing
9. Use bullet points with simple characters like • or -
10. Preserve technical terms and proper nouns exactly as they appear
11. Keep the output concise but complete"""
                },
                {
                    "role": "user",
                    "content": f"Please clean up and format this OCR text with proper line breaks and spacing (no markdown):\n\n{raw_text}"
                }
            ]
            
            response = self.client.chat.completions.create(
                model="llama-3.3-70b-versatile",  # Using fast model for formatting
                messages=messages,
                max_tokens=2000,
                temperature=0.3
            )
            
            formatted_text = response.choices[0].message.content.strip()
            logger.info("Successfully formatted OCR text with AI")
            return formatted_text
            
        except Exception as e:
            logger.warning(f"AI formatting failed, using original text: {e}")
            return raw_text
    
    def _extract_frames_from_video(self, video_path: str, max_frames: int = 10) -> List[str]:
        """Extract frames from video for OCR processing"""
        try:
            cap = cv2.VideoCapture(video_path)
            if not cap.isOpened():
                raise ValueError("Could not open video file")
            
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            duration = total_frames / fps if fps > 0 else 0
            
            # Extract frames at regular intervals
            frame_interval = max(1, total_frames // max_frames)
            frame_paths = []
            
            frame_count = 0
            extracted_count = 0
            
            while cap.isOpened() and extracted_count < max_frames:
                ret, frame = cap.read()
                if not ret:
                    break
                
                if frame_count % frame_interval == 0:
                    # Save frame as temporary image
                    temp_dir = tempfile.gettempdir()
                    frame_path = os.path.join(temp_dir, f"video_frame_{extracted_count}.jpg")
                    cv2.imwrite(frame_path, frame)
                    frame_paths.append(frame_path)
                    extracted_count += 1
                
                frame_count += 1
            
            cap.release()
            logger.info(f"Extracted {len(frame_paths)} frames from video")
            return frame_paths
            
        except Exception as e:
            logger.error(f"Error extracting frames from video: {e}")
            return []
    
    async def extract_text_from_video(self, video_path: str) -> str:
        """Extract text from video by analyzing key frames"""
        try:
            # Extract frames from video
            frame_paths = self._extract_frames_from_video(video_path)
            
            if not frame_paths:
                return "No frames could be extracted from video"
            
            all_text = []
            unique_texts = set()
            
            # Process each frame
            for i, frame_path in enumerate(frame_paths):
                try:
                    frame_text = await self.extract_text_from_image(frame_path)
                    
                    # Only add if text is meaningful and not duplicate
                    if (frame_text and 
                        len(frame_text.strip()) > 10 and 
                        "No readable text found" not in frame_text and
                        "OCR extraction failed" not in frame_text):
                        
                        # Simple deduplication
                        text_hash = hash(frame_text.strip().lower())
                        if text_hash not in unique_texts:
                            unique_texts.add(text_hash)
                            all_text.append(f"--- Frame {i+1} ---\n{frame_text}\n")
                    
                    # Clean up frame file
                    if os.path.exists(frame_path):
                        os.remove(frame_path)
                        
                except Exception as e:
                    logger.warning(f"Error processing frame {i+1}: {e}")
                    continue
            
            if all_text:
                combined_text = "\n".join(all_text)
                logger.info(f"Successfully extracted text from video: {len(combined_text)} characters")
                return combined_text
            else:
                return "No readable text found in video frames"
                
        except Exception as e:
            logger.error(f"Error extracting text from video: {e}")
            return f"Video OCR extraction failed: {str(e)}"
    
    async def extract_text_from_presentation(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract text from tables
                if hasattr(slide, 'shapes'):
                    for shape in slide.shapes:
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    all_text.append("\n".join(slide_text))
            
            combined_text = "\n\n".join(all_text)
            logger.info(f"Successfully extracted text from presentation: {len(combined_text)} characters")
            return combined_text
            
        except Exception as e:
            logger.error(f"Error extracting text from presentation: {e}")
            return f"Presentation text extraction failed: {str(e)}"
    
    def is_image_file(self, filename: str) -> bool:
        """Check if file is an image"""
        image_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp', '.gif'}
        return os.path.splitext(filename.lower())[1] in image_extensions
    
    def is_video_file(self, filename: str) -> bool:
        """Check if file is a video"""
        video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v'}
        return os.path.splitext(filename.lower())[1] in video_extensions
    
    def is_presentation_file(self, filename: str) -> bool:
        """Check if file is a presentation"""
        presentation_extensions = {'.ppt', '.pptx'}
        return os.path.splitext(filename.lower())[1] in presentation_extensions

# Create global instance
ocr_service = OCRService()